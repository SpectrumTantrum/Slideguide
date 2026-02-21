"""
PowerPoint parser using python-pptx for structured slide extraction.

Extracts slide titles, text frames, tables (as markdown), embedded images,
and speaker notes from PPTX files.
"""

from __future__ import annotations

import io
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches  # noqa: F401

from backend.models.schemas import ParsedDocument, SlideContent
from backend.monitoring.logger import get_logger

logger = get_logger(__name__)


class PptxParser:
    """Parse PPTX files into structured slide content."""

    async def parse(self, file_path: str, upload_id: str) -> ParsedDocument:
        """
        Parse a PPTX file into a ParsedDocument.

        Each slide is extracted with its title, text content, tables,
        images, and speaker notes.
        """
        logger.info("pptx_parse_start", file_path=file_path, upload_id=upload_id)

        prs = Presentation(file_path)
        slides: list[SlideContent] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            content = self._extract_slide(slide, slide_num, upload_id)
            slides.append(content)
            logger.debug(
                "pptx_slide_parsed",
                slide_number=slide_num,
                text_length=len(content.text_content),
                has_images=content.has_images,
            )

        result = ParsedDocument(
            upload_id=upload_id,
            file_type="pptx",
            slides=slides,
            metadata={
                "source_file": Path(file_path).name,
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            },
        )

        logger.info(
            "pptx_parse_complete",
            upload_id=upload_id,
            total_slides=result.total_slides,
        )
        return result

    def _extract_slide(
        self, slide, slide_number: int, upload_id: str
    ) -> SlideContent:
        """Extract all content from a single PPTX slide."""
        title = self._get_title(slide)
        text_parts: list[str] = []
        tables: list[str] = []
        image_paths: list[str] = []
        img_idx = 0

        for shape in slide.shapes:
            # Text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)

            # Tables
            if shape.has_table:
                table_md = self._table_to_markdown(shape.table)
                tables.append(table_md)
                text_parts.append(f"\n{table_md}\n")

            # Images
            if hasattr(shape, "image"):
                try:
                    path = self._save_image(
                        shape.image, upload_id, slide_number, img_idx
                    )
                    image_paths.append(path)
                    img_idx += 1
                except Exception as e:
                    logger.warning(
                        "pptx_image_extract_failed",
                        slide_number=slide_number,
                        error=str(e),
                    )

        # Speaker notes
        speaker_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_parts = []
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    notes_parts.append(text)
            speaker_notes = "\n".join(notes_parts)

        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content="\n".join(text_parts),
            has_images=len(image_paths) > 0,
            image_paths=image_paths,
            speaker_notes=speaker_notes,
            tables=tables,
            metadata={"shape_count": len(slide.shapes)},
        )

    def _get_title(self, slide) -> str | None:
        """Extract slide title from title placeholder or first text shape."""
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()

        # Fallback: first text shape with content
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 200:  # Likely a title, not body text
                    return text
        return None

    def _table_to_markdown(self, table) -> str:
        """Convert a PPTX table to markdown format."""
        rows = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if row_idx == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        return "\n".join(rows)

    def _save_image(
        self, image, upload_id: str, slide_number: int, image_index: int
    ) -> str:
        """Save an embedded image to a temp file."""
        temp_dir = Path(tempfile.gettempdir()) / "slideguide" / upload_id
        temp_dir.mkdir(parents=True, exist_ok=True)

        ext = image.content_type.split("/")[-1] if image.content_type else "png"
        image_path = temp_dir / f"slide_{slide_number}_img_{image_index}.{ext}"

        image_path.write_bytes(image.blob)
        return str(image_path)
